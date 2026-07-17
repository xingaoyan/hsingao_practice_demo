#!/usr/bin/env python3
"""
PPTX 解析器
"""

from pathlib import Path
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.slide import Slide

from .base import BaseParser, ParseResult, PageInfo


class PPTXParser(BaseParser):
    """PPTX 解析器"""
    
    def parse(self, file_path: Path, document_id: str) -> ParseResult:
        """解析 PPTX 文件"""
        self.clear_warnings()
        
        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            self.add_warning(f"无法打开 PPTX 文件: {e}")
            return self._create_empty_result(file_path, document_id)
        
        pages = []
        metadata = self._extract_metadata(prs, file_path)
        
        # 提取每张幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            page_info = self._parse_slide(slide, slide_num)
            pages.append(page_info)
        
        # 提取标题
        title = self._extract_title(prs, file_path)
        
        return ParseResult(
            document_id=document_id,
            title=title,
            source_path=str(file_path),
            file_type="pptx",
            pages=pages,
            metadata=metadata,
            warnings=self.warnings
        )
    
    def _parse_slide(self, slide: Slide, slide_number: int) -> PageInfo:
        """解析单张幻灯片"""
        try:
            texts = []
            tables = []
            images = []
            
            for shape in slide.shapes:
                # 提取文本
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                
                # 提取表格
                if shape.has_table:
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        tables.append(table_text)
                
                # 提取图片信息
                if shape.shape_type == 13:  # 图片类型
                    images.append({
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    })
            
            # 提取备注
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
            
            # 组合所有文本
            all_text = "\n".join(texts)
            if notes_text:
                all_text += f"\n\n备注:\n{notes_text}"
            if tables:
                all_text += "\n\n" + "\n".join(tables)
            
            return PageInfo(
                page_number=slide_number,
                text=all_text,
                images=images,
                tables=tables
            )
            
        except Exception as e:
            self.add_warning(f"解析第 {slide_number} 张幻灯片时出错: {e}")
            return PageInfo(
                page_number=slide_number,
                text="",
                images=[],
                tables=[]
            )
    
    def _extract_table_text(self, table) -> str:
        """提取表格文本"""
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.strip())
                rows.append(cells)
            
            if not rows:
                return ""
            
            # 创建 Markdown 表格
            table_lines = []
            table_lines.append("| " + " | ".join(rows[0]) + " |")
            table_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
            for row in rows[1:]:
                table_lines.append("| " + " | ".join(row) + " |")
            
            return "\n".join(table_lines)
            
        except Exception as e:
            self.add_warning(f"提取表格时出错: {e}")
            return ""
    
    def _extract_metadata(self, prs, file_path: Path) -> Dict[str, Any]:
        """提取演示文稿元数据"""
        metadata = self.get_metadata(file_path)
        
        try:
            # 提取核心属性
            core_properties = prs.core_properties
            if core_properties:
                metadata.update({
                    "title": core_properties.title or "",
                    "author": core_properties.author or "",
                    "subject": core_properties.subject or "",
                    "keywords": core_properties.keywords or "",
                    "comments": core_properties.comments or "",
                    "category": core_properties.category or "",
                    "created": str(core_properties.created) if core_properties.created else "",
                    "modified": str(core_properties.modified) if core_properties.modified else "",
                    "last_modified_by": core_properties.last_modified_by or ""
                })
            
            metadata["slide_count"] = len(prs.slides)
            
        except Exception as e:
            self.add_warning(f"提取元数据时出错: {e}")
        
        return metadata
    
    def _extract_title(self, prs, file_path: Path) -> str:
        """提取演示文稿标题"""
        # 首先尝试从元数据获取
        try:
            if prs.core_properties.title:
                return prs.core_properties.title
        except:
            pass
        
        # 尝试从第一张幻灯片提取
        try:
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                return text[:100]
        except:
            pass
        
        # 使用文件名作为标题
        return file_path.stem
    
    def _create_empty_result(self, file_path: Path, document_id: str) -> ParseResult:
        """创建空结果"""
        return ParseResult(
            document_id=document_id,
            title=file_path.stem,
            source_path=str(file_path),
            file_type="pptx",
            pages=[],
            metadata=self.get_metadata(file_path),
            warnings=self.warnings
        )